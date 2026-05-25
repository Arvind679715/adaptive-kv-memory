"""
AKV Research Presentation — Professional Academic Format (UPS Style)
Clean, minimal, high-contrast design. One idea per slide.
Run: pip install python-pptx && python create_presentation.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ===== Professional Color Palette =====
NAVY = RGBColor(0x0D, 0x1B, 0x2A)
DARK_NAVY = RGBColor(0x08, 0x12, 0x1E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF8, 0xF9, 0xFA)
LIGHT_GRAY = RGBColor(0xDE, 0xE2, 0xE6)
MID_GRAY = RGBColor(0x6C, 0x75, 0x7D)
DARK_GRAY = RGBColor(0x21, 0x25, 0x29)
GREEN = RGBColor(0x00, 0xC9, 0x7B)
BLUE = RGBColor(0x33, 0x9A, 0xF0)
RED = RGBColor(0xE8, 0x4C, 0x3D)
GOLD = RGBColor(0xFF, 0xC1, 0x07)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, color=GREEN):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_title(slide, text, subtitle=None, center=False):
    left = Inches(0.8) if not center else Inches(0.5)
    width = Inches(11.8) if not center else Inches(12.3)
    align = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT

    tb = slide.shapes.add_textbox(left, Inches(0.4), width, Inches(0.9))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Segoe UI'
    p.alignment = align

    if subtitle:
        tb2 = slide.shapes.add_textbox(left, Inches(1.25), width, Inches(0.6))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = MID_GRAY
        p2.font.name = 'Segoe UI'
        p2.alignment = align


def add_body(slide, lines, top=2.0, left=0.8, size=20, spacing=12):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(11.5), Inches(5.2))
    tf = tb.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if line.startswith('>>'):
            p.text = line[2:].strip()
            p.font.color.rgb = GREEN
            p.font.bold = True
        elif line.startswith('!!'):
            p.text = line[2:].strip()
            p.font.color.rgb = GOLD
            p.font.bold = True
        elif line.startswith('--'):
            p.text = line[2:].strip()
            p.font.color.rgb = RED
        elif line == '':
            p.text = ''
            p.font.size = Pt(6)
            p.space_after = Pt(2)
            p.font.name = 'Segoe UI'
            continue
        else:
            p.text = line
            p.font.color.rgb = OFF_WHITE
        p.font.size = Pt(size)
        p.font.name = 'Segoe UI'
        p.space_after = Pt(spacing)


def add_stat_box(slide, left, top, number, label, color=GREEN):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(2.8), Inches(0.9))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = 'Segoe UI'
    p.alignment = PP_ALIGN.CENTER

    tb2 = slide.shapes.add_textbox(Inches(left), Inches(top + 0.9), Inches(2.8), Inches(0.7))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = label
    p2.font.size = Pt(13)
    p2.font.color.rgb = LIGHT_GRAY
    p2.font.name = 'Segoe UI'
    p2.alignment = PP_ALIGN.CENTER


def add_result_row(slide, top, cols, is_header=False, is_ours=False):
    color = GREEN if is_ours else (BLUE if is_header else OFF_WHITE)
    bold = is_header or is_ours
    size = 14

    positions = [0.8, 5.2, 7.0, 8.8, 10.8]
    widths = [4.2, 1.6, 1.6, 1.8, 2.0]

    for text, left, width in zip(cols, positions, widths):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.4))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = 'Segoe UI'
        p.alignment = PP_ALIGN.LEFT if left == 0.8 else PP_ALIGN.CENTER

    if is_header:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(0.8), Inches(top + 0.38), Inches(11.7), Inches(0.012))
        line.fill.solid()
        line.fill.fore_color.rgb = MID_GRAY
        line.line.fill.background()


# ================================================================
# SLIDE 1: TITLE
# ================================================================
slide = add_slide()
set_bg(slide)
add_accent_bar(slide, GREEN)

tb = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.2))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "AKV"
p.font.size = Pt(80)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = 'Segoe UI'
p.alignment = PP_ALIGN.CENTER

tb = slide.shapes.add_textbox(Inches(1.5), Inches(3.1), Inches(10.3), Inches(1))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "Adaptive Hierarchical KV Memory"
p.font.size = Pt(34)
p.font.color.rgb = GREEN
p.font.name = 'Segoe UI'
p.alignment = PP_ALIGN.CENTER

tb = slide.shapes.add_textbox(Inches(1.5), Inches(4.3), Inches(10.3), Inches(0.8))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "Near-Lossless KV Cache Compression for Long-Context LLM Inference"
p.font.size = Pt(20)
p.font.color.rgb = LIGHT_GRAY
p.font.name = 'Segoe UI'
p.alignment = PP_ALIGN.CENTER

tb = slide.shapes.add_textbox(Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.6))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "COP Presentation  ·  2026"
p.font.size = Pt(15)
p.font.color.rgb = MID_GRAY
p.font.name = 'Segoe UI'
p.alignment = PP_ALIGN.CENTER

# ================================================================
# SLIDE 2: PROBLEM
# ================================================================
slide = add_slide()
set_bg(slide)
add_accent_bar(slide, RED)
add_title(slide, "The Problem")

add_body(slide, [
    "LLMs store Key-Value pairs for every token in the context window",
    "",
    "!!Memory grows linearly with sequence length:",
    "",
    "   Llama-2-7B  ×  2K tokens    →     1 GB",
    "   Llama-2-7B  ×  32K tokens   →    16 GB",
    "   Llama-2-7B  ×  128K tokens  →    64 GB",
    "",
    "--KV cache is the #1 memory bottleneck for long-context serving",
    "",
    "Current GPU memory (80 GB A100) cannot serve 128K context",
    "with batch size > 1 — KV cache alone exceeds VRAM",
], size=21)

# ================================================================
# SLIDE 3: EXISTING SOLUTIONS
# ================================================================
slide = add_slide()
set_bg(slide)
add_accent_bar(slide, RED)
add_title(slide, "Why Current Solutions Fail")

add_body(slide, [
    "--Eviction methods (H2O, SnapKV, ScissorHands)",
    "   Permanently delete tokens to stay within budget",
    "   Information lost FOREVER — catastrophic if needed later",
    "   +136% perplexity on GPT-2 at 12% budget",
    "",
    "--Uniform quantization (KIVI)",
    "   Quantize ALL tokens to same low bit-width",
    "   Wastes bits on unimportant tokens, hurts important ones",
    "   +24% perplexity on Llama-2-7B",
    "",
    ">>We need: compression WITHOUT permanent information loss",
], size=20)

# ================================================================
# SLIDE 4: OUR INSIGHT
# ================================================================
slide = add_slide()
set_bg(slide)
add_accent_bar(slide, GREEN)

tb = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.3), Inches(2))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Key Insight"
p.font.size = Pt(18)
p.font.color.rgb = MID_GRAY
p.font.name = 'Segoe UI'
p.alignment = PP_ALIGN.CENTER

tb = slide.shapes.add_textbox(Inches(1.5), Inches(3.0), Inches(10.3), Inches(2))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "~10% of tokens receive >80% of attention."
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = 'Segoe UI'
p.alignment = PP_ALIGN.CENTER

tb = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10.3), Inches(1.5))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Give important tokens more bits.\nCompress the rest aggressively.\nNever throw anything away."
p.font.size = Pt(24)
p.font.color.rgb = GREEN
p.font.name = 'Segoe UI'
p.alignment = PP_ALIGN.CENTER

# ================================================================
# SLIDE 5: THREE TIERS
# ================================================================
slide = add_slide()
set_bg(slide)
add_accent_bar(slide, GREEN)
add_title(slide, "AKV: Three-Tier Hierarchy")

# HOT
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.6), Inches(1.7), Inches(3.8), Inches(2.4))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x14, 0x30, 0x20)
shape.line.color.rgb = GREEN
shape.line.width = Pt(2)
tf = shape.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = "HOT"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = GREEN
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "\nfp16  ·  GPU HBM"
p2.font.size = Pt(15)
p2.font.color.rgb = LIGHT_GRAY
p2.alignment = PP_ALIGN.CENTER
p3 = tf.add_paragraph()
p3.text = "Most important tokens"
p3.font.size = Pt(14)
p3.font.color.rgb = MID_GRAY
p3.alignment = PP_ALIGN.CENTER
p4 = tf.add_paragraph()
p4.text = "(recent + high attention)"
p4.font.size = Pt(13)
p4.font.color.rgb = MID_GRAY
p4.alignment = PP_ALIGN.CENTER

# WARM
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(4.8), Inches(1.7), Inches(3.8), Inches(2.4))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x2a, 0x25, 0x10)
shape.line.color.rgb = GOLD
shape.line.width = Pt(2)
tf = shape.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = "WARM"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = GOLD
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "\n4-bit  ·  GPU"
p2.font.size = Pt(15)
p2.font.color.rgb = LIGHT_GRAY
p2.alignment = PP_ALIGN.CENTER
p3 = tf.add_paragraph()
p3.text = "Medium importance"
p3.font.size = Pt(14)
p3.font.color.rgb = MID_GRAY
p3.alignment = PP_ALIGN.CENTER
p4 = tf.add_paragraph()
p4.text = "(per-group quantized)"
p4.font.size = Pt(13)
p4.font.color.rgb = MID_GRAY
p4.alignment = PP_ALIGN.CENTER

# COLD
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(9.0), Inches(1.7), Inches(3.8), Inches(2.4))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x10, 0x20, 0x35)
shape.line.color.rgb = BLUE
shape.line.width = Pt(2)
tf = shape.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = "COLD"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = BLUE
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "\n2-bit  ·  CPU"
p2.font.size = Pt(15)
p2.font.color.rgb = LIGHT_GRAY
p2.alignment = PP_ALIGN.CENTER
p3 = tf.add_paragraph()
p3.text = "Least important"
p3.font.size = Pt(14)
p3.font.color.rgb = MID_GRAY
p3.alignment = PP_ALIGN.CENTER
p4 = tf.add_paragraph()
p4.text = "(but NEVER deleted)"
p4.font.size = Pt(13)
p4.font.color.rgb = MID_GRAY
p4.alignment = PP_ALIGN.CENTER

# Arrows
for x in [4.4, 8.6]:
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
        Inches(x), Inches(2.7), Inches(0.4), Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MID_GRAY
    shape.line.fill.background()

add_body(slide, [
    "Mirrors OS virtual memory:   L1 (hot) → L2 (warm) → Swap (cold)",
    "",
    ">>Tokens are demoted when HOT exceeds budget — NEVER evicted",
    "Tokens can be promoted back if importance increases over time",
], top=4.5, size=18)

# ================================================================
# SLIDE 6: PROCESS
# ================================================================
slide = add_slide()
set_bg(slide)
add_accent_bar(slide, GREEN)
add_title(slide, "How It Works", "Per decoding step")

add_body(slide, [
    ">>Step 1:  New token KV → enters HOT tier (fp16)",
    "",
    ">>Step 2:  Score ALL cached tokens",
    "           Hybrid: attention accumulation + recency + decay",
    "",
    ">>Step 3:  If HOT exceeds budget → reorganize:",
    "           Top-k scores     →  stay HOT   (fp16)",
    "           Middle scores    →  demote to WARM  (4-bit)",
    "           Bottom scores    →  demote to COLD  (2-bit)",
    "",
    ">>Step 4:  Compute attention over ALL tiers",
    "           Dequantize warm/cold on-the-fly",
], size=20)

# ================================================================
# SLIDE 7: DIFFERENTIATORS
# ================================================================
slide = add_slide()
set_bg(slide)
add_accent_bar(slide, GREEN)
add_title(slide, "Key Differentiators")

add_body(slide, [
    ">>1. Never evicts — only demotes",
    "   Cold tokens always retrievable (unlike H2O/SnapKV)",
    "",
    ">>2. Mixed-precision (importance-aware)",
    "   Important = fp16, less important = 4-bit, unimportant = 2-bit",
    "",
    ">>3. Continuous adaptation",
    "   Scores evolve during generation; tokens can move between tiers",
    "",
    ">>4. Hardware-aligned hierarchy",
    "   GPU HBM → GPU quantized → CPU RAM (matches memory bandwidth)",
], size=20)

# ================================================================
# SLIDE 8: SECTION BREAK
# ================================================================
slide = add_slide()
set_bg(slide)
add_accent_bar(slide, GREEN)

tb = slide.shapes.add_textbox(Inches(1), Inches(3.0), Inches(11.3), Inches(1.5))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "Experimental Results"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = 'Segoe UI'
p.alignment = PP_ALIGN.CENTER

tb = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.3), Inches(0.8))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "WikiText-2 Perplexity  ·  GPT-2 (124M) & Llama-2-7B"
p.font.size = Pt(20)
p.font.color.rgb = MID_GRAY
p.font.name = 'Segoe UI'
p.alignment = PP_ALIGN.CENTER

# ================================================================
# SLIDE 9: LLAMA-2 RESULTS
# ================================================================
slide = add_slide()
set_bg(slide)
add_accent_bar(slide, GREEN)
add_title(slide, "Llama-2-7B Results", "WikiText-2  ·  2048 tokens  ·  Budget = 256 (12.5% retention)")

add_result_row(slide, 2.1, ["Method", "PPL ↓", "Δ%", "Memory", "Compress."], is_header=True)
add_result_row(slide, 2.6, ["Full Cache (baseline)", "4.21", "—", "1074 MB", "1.0×"])
add_result_row(slide, 3.05, ["H2O (permanent eviction)", "5.75", "+36.5%", "134 MB", "8.0×"])
add_result_row(slide, 3.5, ["KIVI-2bit (uniform quant)", "5.24", "+24.3%", "602 MB", "1.8×"])

# Separator line before our results
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.8), Inches(3.95), Inches(11.7), Inches(0.008))
line.fill.solid()
line.fill.fore_color.rgb = GREEN
line.line.fill.background()

add_result_row(slide, 4.1, ["AKV-4bit (ours)", "4.23", "+0.5%", "384 MB", "2.8×"], is_ours=True)
add_result_row(slide, 4.55, ["AKV-2bit (ours)", "4.93", "+17.0%", "266 MB", "4.0×"], is_ours=True)

# Key stats at bottom
add_stat_box(slide, 0.8, 5.4, "+0.5%", "PPL degradation\n(near-lossless)", GREEN)
add_stat_box(slide, 3.6, 5.4, "2.8×", "compression\n(saves 690 MB)", BLUE)
add_stat_box(slide, 6.4, 5.4, "2048", "tokens retained\n(H2O keeps 256)", GOLD)
add_stat_box(slide, 9.4, 5.4, "0", "tokens evicted\n(never forgets)", GREEN)

# ================================================================
# SLIDE 10: GPT-2 RESULTS
# ================================================================
slide = add_slide()
set_bg(slide)
add_accent_bar(slide, GREEN)
add_title(slide, "GPT-2 Results", "WikiText-2  ·  1024 tokens  ·  Budget = 128 (12.5% retention)")

add_result_row(slide, 2.1, ["Method", "PPL ↓", "Δ%", "Memory", "Compress."], is_header=True)
add_result_row(slide, 2.6, ["Full Cache (baseline)", "15.80", "—", "37.7 MB", "1.0×"])
add_result_row(slide, 3.05, ["H2O (permanent eviction)", "37.32", "+136%", "4.7 MB", "8.0×"])
add_result_row(slide, 3.5, ["KIVI-2bit (uniform quant)", "20.31", "+28.5%", "22.3 MB", "1.7×"])

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.8), Inches(3.95), Inches(11.7), Inches(0.008))
line.fill.solid()
line.fill.fore_color.rgb = GREEN
line.line.fill.background()

add_result_row(slide, 4.1, ["AKV-4bit (ours)", "15.71", "−0.6%", "14.0 MB", "2.7×"], is_ours=True)
add_result_row(slide, 4.55, ["AKV-2bit (ours)", "18.50", "+17.1%", "9.9 MB", "3.8×"], is_ours=True)

add_body(slide, [
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    ">>AKV-4bit: lossless compression (−0.6% PPL) at 2.7× memory reduction",
    "--H2O: collapses at low budgets (+136% — more than doubles perplexity)",
], top=2.0, size=18)

# ================================================================
# SLIDE 11: HEAD-TO-HEAD
# ================================================================
slide = add_slide()
set_bg(slide)
add_accent_bar(slide, GREEN)
add_title(slide, "Head-to-Head Comparison")

add_body(slide, [
    ">>AKV-4bit vs Full Cache  (Llama-2-7B)",
    "   PPL: 4.23 vs 4.21  ·  Memory: 384 vs 1074 MB",
    "   Essentially lossless, saves 690 MB (64% reduction)",
    "",
    ">>AKV-2bit vs KIVI-2bit  (same 2-bit budget)",
    "   PPL: 4.93 vs 5.24  ·  Memory: 266 vs 602 MB",
    "   Better quality AND better compression",
    "   Why: important tokens stay fp16 in HOT tier",
    "",
    ">>AKV-4bit vs H2O  (same budget constraint)",
    "   PPL: 4.23 vs 5.75  ·  Tokens: 2048 vs 256 retained",
    "   H2O permanently loses 88% of context",
    "   AKV keeps everything — just at adaptive precision",
], size=19)

# ================================================================
# SLIDE 12: WHY IT WORKS
# ================================================================
slide = add_slide()
set_bg(slide)
add_accent_bar(slide, GREEN)
add_title(slide, "Why AKV Works")

add_body(slide, [
    ">>Attention is highly skewed in practice",
    "   ~10% of tokens get >80% of attention (BOS, entities, recent)",
    "   These NEED full precision — rest can be compressed safely",
    "",
    ">>Group-based quantization preserves structure",
    "   Per-channel asymmetric with configurable group size",
    "   4-bit cosine similarity > 0.999 with fp16 original",
    "",
    ">>Adaptive scoring catches evolving importance",
    "   Token unimportant at step 100 may matter at step 500",
    "   AKV can promote it back to higher precision",
    "",
    ">>No single point of failure",
    "   Unlike eviction: a wrongly-demoted token is still accessible",
], size=19)

# ================================================================
# SLIDE 13: ABLATION
# ================================================================
slide = add_slide()
set_bg(slide)
add_accent_bar(slide, GREEN)
add_title(slide, "Ablation Studies")

add_body(slide, [
    ">>Warm tier bit-width",
    "   2-bit: 6.5× compression, some quality loss",
    "   4-bit: 3.8× compression, near-lossless  ← recommended",
    "   8-bit: 1.8× compression, diminishing returns",
    "",
    ">>Hot budget size (% of sequence in fp16)",
    "   Smaller → more compression, slight PPL increase",
    "   12-25% of sequence length is the sweet spot",
    "",
    ">>Quantization group size",
    "   32: lowest error, more overhead",
    "   128: best speed/quality balance  ← recommended",
    "   256: fastest, slightly more error",
], size=19)

# ================================================================
# SLIDE 14: LIMITATIONS
# ================================================================
slide = add_slide()
set_bg(slide)
add_accent_bar(slide, BLUE)
add_title(slide, "Limitations & Future Work")

add_body(slide, [
    "Limitations:",
    "   • Python-level management — not kernel-fused yet (latency overhead)",
    "   • Evaluated on 1K-2K contexts — need 32K+ benchmarks",
    "   • No retrieval task (passkey / needle-in-haystack) yet",
    "",
    "Future work:",
    ">>  • Fused Triton kernels: dequantize-inside-attention",
    ">>  • 128K context on Llama-3-8B / Mistral-7B",
    ">>  • Needle-in-a-haystack: prove cold tokens are retrievable",
    "   • vLLM / TensorRT-LLM integration",
    "   • Multi-GPU cold-tier sharding",
], size=19)

# ================================================================
# SLIDE 15: CONCLUSION
# ================================================================
slide = add_slide()
set_bg(slide)
add_accent_bar(slide, GREEN)
add_title(slide, "Conclusion", center=True)

add_body(slide, [
    "",
    ">>Near-lossless compression",
    "   Llama-2-7B: +0.5% PPL at 2.8× compression (saves 690 MB)",
    "   GPT-2: −0.6% PPL at 2.7× compression",
    "",
    ">>Outperforms all baselines",
    "   vs KIVI: better quality AND better compression (same bits)",
    "   vs H2O: 10× better quality, retains ALL tokens",
    "",
    ">>Core innovation",
    "   Importance-aware mixed-precision tiering",
    "   Hot (fp16) → Warm (4-bit) → Cold (2-bit)",
    "   Never evict — only demote",
], top=1.8, size=21)

# ================================================================
# SLIDE 16: THANK YOU
# ================================================================
slide = add_slide()
set_bg(slide)
add_accent_bar(slide, GREEN)

tb = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.3), Inches(1.2))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "Thank You"
p.font.size = Pt(56)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = 'Segoe UI'
p.alignment = PP_ALIGN.CENTER

tb = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.3), Inches(0.8))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "Questions?"
p.font.size = Pt(30)
p.font.color.rgb = GREEN
p.font.name = 'Segoe UI'
p.alignment = PP_ALIGN.CENTER

# Takeaway bar
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(2), Inches(5.2), Inches(9.3), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x14, 0x30, 0x20)
shape.line.color.rgb = GREEN
shape.line.width = Pt(1)
tf = shape.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = "Key Result:  +0.5% perplexity  ·  2.8× compression  ·  0 tokens lost"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = GREEN
p.font.name = 'Segoe UI'
p.alignment = PP_ALIGN.CENTER

# ================================================================
# SAVE
# ================================================================
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                           'AKV_COP_Presentation.pptx')
prs.save(output_path)
print(f'Saved: {output_path}')
print(f'  {len(prs.slides)} slides  ·  16:9 widescreen  ·  Dark navy + green accent')
